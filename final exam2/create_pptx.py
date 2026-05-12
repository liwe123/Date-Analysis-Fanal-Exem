# -*- coding: utf-8 -*-
"""create_pptx.py - 生成大数据学期项目 RAG 答辩演示文稿（全中文精美设计）"""

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# ===== 颜色定义 =====
D  = RGBColor(0x0D,0x23,0x3F)  # 深蓝黑
B  = RGBColor(0x1A,0x56,0xDB)  # 主蓝
T  = RGBColor(0x14,0x8A,0x8A)  # 青绿
G  = RGBColor(0x10,0xB9,0x81)  # 绿
A  = RGBColor(0x06,0xD6,0xA0)  # 亮青
O  = RGBColor(0xF5,0x9E,0x0B)  # 橙
R  = RGBColor(0xEF,0x44,0x44)  # 红
P  = RGBColor(0x8B,0x5C,0xF6)  # 紫
CA = RGBColor(0xF8,0xFA,0xFC)  # 卡片背景
TX = RGBColor(0x1E,0x29,0x3B)  # 正文
GR = RGBColor(0x6B,0x72,0x80)  # 灰色
LG = RGBColor(0x94,0xA3,0xB8)  # 浅灰
WH = RGBColor(0xFF,0xFF,0xFF)  # 白色

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width; H = prs.slide_height

# ===== 工具函数 =====
def bg(sl,c):
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = c

def re(sl,l,t,w,h,f=None):
    sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h); sh.line.fill.background()
    if f: sh.fill.solid(); sh.fill.fore_color.rgb = f
    else: sh.fill.background()
    return sh

def ov(sl,l,t,s,f=None):
    sh = sl.shapes.add_shape(MSO_SHAPE.OVAL,l,t,s,s); sh.line.fill.background()
    if f: sh.fill.solid(); sh.fill.fore_color.rgb = f
    return sh

def tx(sl,l,t,w,h,txt='',sz=14,c=TX,b=False,a=PP_ALIGN.LEFT):
    box = sl.shapes.add_textbox(l,t,w,h); tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
    p.font.name = 'Microsoft YaHei'; p.alignment = a
    return box,tf

def ap(tf,txt,sz=14,c=TX,b=False,a=PP_ALIGN.LEFT,be=0,af=4):
    p = tf.add_paragraph(); p.text = txt
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
    p.font.name = 'Microsoft YaHei'; p.alignment = a
    p.space_before = Pt(be); p.space_after = Pt(af)
    return p

def hdr(sl,tit,sub=''):
    re(sl,0,0,W,Inches(1.0),D); re(sl,0,Inches(0.95),W,Inches(0.05),A)
    tx(sl,Inches(0.6),Inches(0.12),Inches(10),Inches(0.5),tit,26,WH,True)
    if sub: tx(sl,Inches(0.6),Inches(0.55),Inches(10),Inches(0.35),sub,14,LG)

def card(sl,l,t,w,h,tit,items,color=B):
    re(sl,l,t,w,h,CA)
    re(sl,l+Inches(0.02),t+Inches(0.02),w-Inches(0.04),Inches(0.04),color)
    tx(sl,l+Inches(0.2),t+Inches(0.15),w-Inches(0.4),Inches(0.3),tit,16,color,True)
    for i,item in enumerate(items):
        tx(sl,l+Inches(0.2),t+Inches(0.55)+Inches(0.28)*i,w-Inches(0.4),Inches(0.26),
           '  '+item,12,TX)

# ====================================================================
# 第 1 页 · 封面
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,D)
ov(sl,Inches(-0.5),Inches(-0.5),Inches(2.0),B)
ov(sl,Inches(11.5),Inches(-1.0),Inches(3.0),T)
ov(sl,Inches(10.5),Inches(5.5),Inches(4.0),B)
ov(sl,Inches(-1.0),Inches(6.0),Inches(2.5),T)
re(sl,0,Inches(2.4),W,Inches(2.8),B)
tx(sl,Inches(0.8),Inches(2.6),Inches(11.5),Inches(0.9),
   '大数据学期项目计划二',42,WH,True,PP_ALIGN.CENTER)
tx(sl,Inches(0.8),Inches(3.5),Inches(11.5),Inches(0.8),
   'RAG 检索增强生成助手',30,WH,False,PP_ALIGN.CENTER)
tx(sl,Inches(1.5),Inches(5.6),Inches(10),Inches(0.4),
   'ChromaDB  |  DeepSeek V4  |  Qwen3-Embedding  |  Streamlit',15,LG,False,PP_ALIGN.CENTER)
tx(sl,Inches(1.5),Inches(6.1),Inches(10),Inches(0.4),
   '2025-2026 学年  |  小组答辩',14,LG,False,PP_ALIGN.CENTER)

# ====================================================================
# 第 2 页 · 目录
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,WH)
re(sl,0,0,W,Inches(1.0),D); re(sl,0,Inches(0.95),W,Inches(0.05),A)
tx(sl,Inches(0.6),Inches(0.15),Inches(10),Inches(0.7),'目  录',28,WH,True)
toc = [('01','项目背景与目标',B),('02','系统架构总览',T),
       ('03','数据流水线详解',G),('04','检索与答案生成',O),
       ('05','代码质量与测试',P),('06','项目运行数据',B),
       ('07','项目创新亮点',T),('08','局限与改进方向',G),('09','演示环节',O)]
for i,(nu,ti,co) in enumerate(toc):
    col = i//5; off = i%5
    x = Inches(1.2)+Inches(6.0)*col; y = Inches(1.4)+Inches(1.05)*off
    ov(sl,x,y+Inches(0.05),Inches(0.5),co)
    tx(sl,x+Inches(0.05),y+Inches(0.1),Inches(0.4),Inches(0.4),nu,14,WH,True,PP_ALIGN.CENTER)
    tx(sl,x+Inches(0.7),y+Inches(0.1),Inches(4.5),Inches(0.4),ti,17,TX)

# ====================================================================
# 第 3 页 · 背景与目标
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,WH)
hdr(sl,'一、项目背景与目标')
re(sl,Inches(0.6),Inches(1.5),Inches(5.8),Inches(5.5),CA)
ov(sl,Inches(1.0),Inches(1.7),Inches(0.5),R)
tx(sl,Inches(1.0),Inches(1.72),Inches(0.5),Inches(0.45),'!',18,WH,True,PP_ALIGN.CENTER)
tx(sl,Inches(1.7),Inches(1.7),Inches(4.5),Inches(0.4),'企业痛点',18,R,True)
for i,it in enumerate(['知识分散在各类非结构化文档中','关键词搜索无法理解语义',
    '团队重复回答相同问题效率低','新员工上手慢知识靠口口相传']):
    tx(sl,Inches(1.0),Inches(2.3)+Inches(0.55)*i,Inches(5.2),Inches(0.5),it,13,TX)
re(sl,Inches(6.9),Inches(1.5),Inches(5.8),Inches(5.5),CA)
ov(sl,Inches(7.3),Inches(1.7),Inches(0.5),G)
tx(sl,Inches(7.3),Inches(1.72),Inches(0.5),Inches(0.45),'!',18,WH,True,PP_ALIGN.CENTER)
tx(sl,Inches(8.0),Inches(1.7),Inches(4.5),Inches(0.4),'项目目标',18,G,True)
for i,it in enumerate(['构建完整 RAG 问答流水线','全链路自动化处理','语义检索加元数据混合搜索','带来源追溯防幻觉','CLI 加 Streamlit 双模式']):
    tx(sl,Inches(7.3),Inches(2.3)+Inches(0.55)*i,Inches(5.2),Inches(0.5),it,13,TX)

# ====================================================================
# 第 4 页 · 系统架构
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,WH)
hdr(sl,'二、系统架构总览','七层流水线：从原始文档到带来源追溯的智能问答')
layers = [('1\n数据摄取',B),('2\n文本预处理',T),('3\n向量嵌入',G),
          ('4\nChromaDB\n存储',A),('5\n查询解析',O),('6\n混合检索',R),('7\n答案生成',P)]
bw = Inches(1.5); bh = Inches(1.3); gap = Inches(0.12)
tot = 7*bw+6*gap; sx = (W-tot)//2
for i,(lb,co) in enumerate(layers):
    x = sx+i*(bw+gap); sh = re(sl,x,Inches(2.3),bw,bh,co); sh.text_frame.word_wrap = True
    ls = lb.split('\n'); sh.text_frame.paragraphs[0].text = ls[0]
    for p in sh.text_frame.paragraphs:
        p.font.size = Pt(13); p.font.color.rgb = WH; p.font.bold = True
        p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER
    if len(ls)>1:
        p2 = sh.text_frame.add_paragraph(); p2.text = ls[1]; p2.font.size = Pt(10)
        p2.font.color.rgb = WH; p2.font.name = 'Microsoft YaHei'; p2.alignment = PP_ALIGN.CENTER
    if i<len(layers)-1:
        tx(sl,x+bw-Inches(0.05),Inches(2.75),Inches(0.25),Inches(0.3),'>',20,LG,True,PP_ALIGN.CENTER)
tx(sl,Inches(0.6),Inches(4.0),Inches(12),Inches(0.35),
   '输入：data/raw/ -> 98 篇文档（46 篇课程资料 + 52 篇 Wikipedia 词条）',14,GR)
tx(sl,Inches(0.6),Inches(4.4),Inches(12),Inches(0.35),'输出：CLI 命令行  |  Streamlit Web 界面',14,GR)
_,tf = tx(sl,Inches(0.6),Inches(5.6),Inches(12),Inches(0.3),'',12)
ap(tf,'数据流：读取 -> 清洗分块 -> LLM 元数据 -> Qwen3 嵌入 -> ChromaDB 存储 -> 解析 -> 混合检索 -> DeepSeek 生成',12,GR)

# ====================================================================
# 第 5 页 · 数据流水线
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,WH)
hdr(sl,'三、数据流水线详解','文档摄取与语料自动采集')
card(sl,Inches(0.6),Inches(1.5),Inches(5.8),Inches(5.5),'文档加载 ingest.py',
   ['* 支持 .md / .txt / .pdf 格式','* 递归遍历 data/raw/ 目录','* 解析 YAML Front-Matter',
    '* 返回结构化文档对象','* UTF-8 容错读取'],B)
card(sl,Inches(6.9),Inches(1.5),Inches(5.8),Inches(5.5),'语料采集 collect_corpus.py',
   ['* Wikipedia REST API 免费调用','* 58 个话题成功采集 52 个','* 优先中文失败回退英文',
    '* 限速保护 2 秒间隔 3 次重试','* 429 自动等待 404 直接跳过','* 输出带 Front-Matter 的 MD 文件'],T)

# ====================================================================
# 第 6 页 · 清洗与分块
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,WH)
hdr(sl,'三、文本清洗与语义分块')
tx(sl,Inches(0.6),Inches(1.4),Inches(12),Inches(0.4),'文本清洗四步流水线 clean_text',18,B,True)
for i,(t,d) in enumerate([('第一步：移除 HTML 标签','正则 <[^>]+> 替换为空格'),
    ('第二步：解码 HTML 实体','&amp; -> &   &lt; -> <   &nbsp; -> 空格'),
    ('第三步：过滤控制字符','保留换行符和制表符，移除不可打印字符'),
    ('第四步：空白规范化','合并连续空格，压缩多余换行')]):
    y = Inches(2.0)+Inches(0.55)*i; re(sl,Inches(0.8),y,Inches(11.5),Inches(0.45),CA)
    tx(sl,Inches(1.0),y+Inches(0.05),Inches(2.5),Inches(0.35),t,14,B,True)
    tx(sl,Inches(3.5),y+Inches(0.05),Inches(8.5),Inches(0.35),d,13,TX)
tx(sl,Inches(0.6),Inches(4.2),Inches(12),Inches(0.4),'语义分块四层优先级 chunk_text',18,T,True)
for i,(t,d) in enumerate([('第一层：段落边界','按空行切分，保持段落语义完整'),
    ('第二层：句子边界','按句末标点切分，保留完整句子'),
    ('第三层：贪心合并','合并到 chunk_size=700 字符，重叠 120 字符'),
    ('第四层：滑窗切割','超长文本以 step=chunk_size-overlap 切割')]):
    y = Inches(4.7)+Inches(0.5)*i; re(sl,Inches(0.8),y,Inches(0.04),Inches(0.4),T)
    tx(sl,Inches(1.1),y,Inches(3.0),Inches(0.4),t,13,T,True); tx(sl,Inches(4.1),y,Inches(8.0),Inches(0.4),d,12,TX)

# ====================================================================
# 第 7 页 · 元数据
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,WH)
hdr(sl,'三、LLM 元数据提取与 JSON 安全解析')
card(sl,Inches(0.6),Inches(1.5),Inches(5.8),Inches(5.0),'LLM 元数据提取',
   ['* 发送前 1200 字符给 DeepSeek V4','* 返回 JSON：作者/年份/分类/语言/摘要',
    '* Front-Matter 优先覆盖 LLM 结果','* 文件名后缀猜测兜底','* 98 篇文档成本小于 0.5 元'],G)
card(sl,Inches(6.9),Inches(1.5),Inches(5.8),Inches(5.0),'JSON 安全解析',
   ['* 不因 LLM 格式问题中断流水线','* 第一层：标准 json.loads','* 第二层：截断到完整大括号对象',
    '* 第三层：移除注释和尾随逗号','* 第四层：逐字段正则提取'],P)
re(sl,Inches(0.6),Inches(6.2),Inches(12.1),Inches(0.8),CA)
re(sl,Inches(0.6),Inches(6.2),Inches(0.06),Inches(0.8),A)
tx(sl,Inches(0.9),Inches(6.25),Inches(11.5),Inches(0.3),
   '合并策略：Front-Matter 人工标注 > LLM 提取 > 文件名猜测 -> 最终写入 ChromaDB',15,A,True)

# ====================================================================
# 第 8 页 · 向量存储
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,WH)
hdr(sl,'四、检索与答案生成','向量存储与混合检索')
card(sl,Inches(0.6),Inches(1.5),Inches(5.8),Inches(5.0),'向量存储配置',
   ['* 数据库：ChromaDB 本地持久化','* 距离度量：Cosine 余弦距离',
    '* 索引算法：HNSW','* 写入策略：批量 upsert 每次 64 条',
    '* 双模式嵌入：本地/远程自由切换','* 存储路径：vector_store/'],B)
card(sl,Inches(6.9),Inches(1.5),Inches(5.8),Inches(5.0),'混合检索三大能力',
   ['* 1. 纯语义搜索：向量相似度 Top-K','* 2. 语义加 Where 过滤：按年份/分类筛选',
    '* 3. 最大距离阈值：过滤低质量结果','* 4. Where 失败自动回退纯语义搜索'],O)
re(sl,Inches(0.6),Inches(6.2),Inches(12.1),Inches(0.8),CA)
re(sl,Inches(0.6),Inches(6.2),Inches(0.06),Inches(0.8),O)
tx(sl,Inches(0.9),Inches(6.25),Inches(11.5),Inches(0.3),
   '降级保障：Where 过滤失败 -> 自动移除过滤条件 -> 纯语义搜索回退 -> 用户永不见空白',15,O,True)

# ====================================================================
# 第 9 页 · 本地嵌入
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,WH)
hdr(sl,'四、本地嵌入模型（GPU 加速）')
re(sl,Inches(0.6),Inches(1.5),Inches(12.1),Inches(2.5),CA)
re(sl,Inches(0.6),Inches(1.5),Inches(0.06),Inches(2.5),G)
tx(sl,Inches(1.0),Inches(1.7),Inches(11),Inches(0.4),'Qwen3-Embedding-0.6B',22,G,True)
for i,f in enumerate(['* 参数量 0.6B，输出 1024 维向量，支持 32K 上下文',
    '* 中文 MTEB 66.33 同级最优','* 自动检测 CUDA 启用 GPU 加速',
    '* 编码 130 多个文本块仅需 10-20 秒（GPU 加速）','* 通过环境变量可一键切换任意 HuggingFace 模型']):
    tx(sl,Inches(1.0),Inches(2.2)+Inches(0.35)*i,Inches(11),Inches(0.3),f,13,TX)
tx(sl,Inches(0.6),Inches(4.3),Inches(12),Inches(0.4),'双模式嵌入切换',18,B,True)
cx=[Inches(0.8),Inches(3.8),Inches(6.5),Inches(8.5),Inches(10.5)]
cw=[Inches(2.8),Inches(2.5),Inches(1.8),Inches(1.8),Inches(2.0)]
for i,(x,w) in enumerate(zip(cx,cw)): re(sl,x,Inches(4.8),w,Inches(0.45),D)
for i,(x,w,h) in enumerate(zip(cx,cw,['模式','模型','维度','成本','速度'])):
    tx(sl,x+Inches(0.1),Inches(4.85),w-Inches(0.2),Inches(0.35),h,13,WH,True,PP_ALIGN.CENTER)
tr=[['本地（当前）','Qwen3-0.6B','1024','0 元','GPU 10-20 秒'],
    ['远程 API','text-embed-3-small','1536','约 0.01 元','API 延迟']]
for ri,row in enumerate(tr):
    y = Inches(5.3)+Inches(0.45)*ri; bg2 = CA if ri%2==0 else WH
    for ci,(x,w) in enumerate(zip(cx,cw)):
        re(sl,x,y,w,Inches(0.45),bg2)
        tx(sl,x+Inches(0.1),y+Inches(0.05),w-Inches(0.2),Inches(0.35),row[ci],12,TX,ci==0,PP_ALIGN.CENTER)

# ====================================================================
# 第 10 页 · 查询解析
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,WH)
hdr(sl,'四、查询意图解析','将自然语言自动转为结构化搜索参数')
tx(sl,Inches(0.8),Inches(1.4),Inches(11),Inches(0.4),
   'LLM 自动将用户问题拆解为搜索词加过滤条件，用户无需学习搜索语法。',14,GR)
ex=[('用户提问：去年期末怎么考？','解析结果：{"search_query": "期末考试 形式", "filters": {"year": 2025, "category": "notice"}}'),
    ('用户提问：2024 年通知讲了啥？','解析结果：{"search_query": "通知", "filters": {"year": 2024, "category": "notice"}}'),
    ('用户提问：什么是向量数据库？','解析结果：{"search_query": "向量数据库", "filters": null}'),
    ('用户提问：张三写的案例分析？','解析结果：{"search_query": "案例分析", "filters": {"author": "张三"}}')]
for i,(q,a) in enumerate(ex):
    y = Inches(2.0)+Inches(1.0)*i; re(sl,Inches(0.8),y,Inches(11.7),Inches(0.85),CA)
    re(sl,Inches(0.8),y,Inches(0.06),Inches(0.85),B)
    tx(sl,Inches(1.1),y+Inches(0.05),Inches(11),Inches(0.3),q,14,B,True)
    tx(sl,Inches(1.1),y+Inches(0.4),Inches(11),Inches(0.35),a,12,GR)
re(sl,Inches(0.8),Inches(6.3),Inches(11.7),Inches(0.7),CA)
re(sl,Inches(0.8),Inches(6.3),Inches(0.06),Inches(0.7),O)
tx(sl,Inches(1.1),Inches(6.35),Inches(11),Inches(0.5),
   '降级保护：JSON 解析异常或 API 调用失败时，自动使用原问题全文搜索，系统永不卡死。',13,O,True)

# ====================================================================
# 第 11 页 · 防幻觉
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,WH)
hdr(sl,'四、答案生成与防幻觉机制')
for i,(nu,ti,de,co) in enumerate([('1','System Prompt 约束','LLM 被要求仅依据\n参考资料回答，找不到\n就说不知道。',G),
    ('2','上下文隔离','只将检索结果喂给\nLLM，不给模型自我\n发挥的空间。',G),
    ('3','引用强制检查','后端自动检查答案中\n是否含 [Source:]，\n缺失则自动补全。',G)]):
    x = Inches(0.6)+Inches(4.2)*i; re(sl,x,Inches(1.5),Inches(3.8),Inches(3.2),CA)
    ov(sl,x+Inches(1.5),Inches(1.7),Inches(0.6),co)
    tx(sl,x+Inches(1.55),Inches(1.72),Inches(0.5),Inches(0.55),nu,16,WH,True,PP_ALIGN.CENTER)
    tx(sl,x+Inches(0.2),Inches(2.5),Inches(3.4),Inches(0.4),ti,16,co,True,PP_ALIGN.CENTER)
    tx(sl,x+Inches(0.2),Inches(2.9),Inches(3.4),Inches(0.8),de,12,TX,False,PP_ALIGN.CENTER)
tx(sl,Inches(0.6),Inches(5.0),Inches(12),Inches(0.4),'效果验证',18,D,True)
re(sl,Inches(0.6),Inches(5.5),Inches(5.8),Inches(1.5),CA)
re(sl,Inches(0.6),Inches(5.5),Inches(0.06),Inches(1.5),G)
tx(sl,Inches(0.9),Inches(5.6),Inches(5.3),Inches(0.3),'知识库有答案',15,G,True)
tx(sl,Inches(0.9),Inches(5.95),Inches(5.3),Inches(0.7),
   '问：课程最后提交截止日期是什么？\n答：答案 + 【来源：xxx.md】',12,TX)
re(sl,Inches(6.9),Inches(5.5),Inches(5.8),Inches(1.5),CA)
re(sl,Inches(6.9),Inches(5.5),Inches(0.06),Inches(1.5),O)
tx(sl,Inches(7.2),Inches(5.6),Inches(5.3),Inches(0.3),'知识库无答案',15,O,True)
tx(sl,Inches(7.2),Inches(5.95),Inches(5.3),Inches(0.7),
   '问：钢琴考级需要准备什么？\n答：资料中没有相关信息。',12,TX)

# ====================================================================
# 第 12 页 · 代码质量
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,WH)
hdr(sl,'五、代码质量与工程实践')
card(sl,Inches(0.6),Inches(1.5),Inches(5.8),Inches(5.5),'模块化架构',
   ['* main.py（118 行）- CLI 入口','* ingest.py（104 行）- 文档加载',
    '* preprocess.py（301 行）- 清洗分块','* embed_store.py（210 行）- 向量存储',
    '* qa.py（71 行）- 答案生成','* query_parser.py（113 行）- 查询解析',
    '* collect_corpus.py（199 行）- 语料采集','* utils.py（79 行）- 工具函数'],B)
card(sl,Inches(6.9),Inches(1.5),Inches(5.8),Inches(5.5),'测试覆盖（45 个以上用例）',
   ['* test_preprocess - 清洗/分块/分类','* test_ingest - Front-Matter 解析',
    '* test_query_parser - 正常/异常/回退','* test_embed_store - 搜索/过滤/删除',
    '* test_qa - 空文档/上下文/引用','* test_integration - 全链路集成',
    '','* 每个公开函数覆盖正常加回退路径'],T)
re(sl,Inches(0.6),Inches(6.2),Inches(12.1),Inches(0.8),CA)
re(sl,Inches(0.6),Inches(6.2),Inches(0.06),Inches(0.8),P)
tx(sl,Inches(0.9),Inches(6.25),Inches(11.5),Inches(0.3),
   '工程规范：类型注解全覆盖 | 模块间单向依赖 | 显式 init_env | OpenAI 单例 | 分层错误处理',15,P,True)

# ====================================================================
# 第 13 页 · 运行数据
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,WH)
hdr(sl,'六、项目运行数据')
re(sl,Inches(0.6),Inches(1.5),Inches(5.8),Inches(4.5),CA)
tx(sl,Inches(0.8),Inches(1.7),Inches(5.4),Inches(0.4),'数据规模',18,B,True)
for i,(la,va,cl) in enumerate([('课程文档','46 篇',B),('Wikipedia 词条','52 篇',T),
    ('文档总数','98 篇',D),('文档分类','7 大类',G),('处理后文本块','约 130-150 个',O)]):
    y = Inches(2.2)+Inches(0.55)*i; tx(sl,Inches(0.9),y,Inches(2.0),Inches(0.35),la,13,GR)
    re(sl,Inches(3.0),y,Inches(0.04),Inches(0.35),cl); tx(sl,Inches(3.2),y,Inches(3.0),Inches(0.35),va,14,cl,True)
re(sl,Inches(6.9),Inches(1.5),Inches(5.8),Inches(4.5),CA)
tx(sl,Inches(7.1),Inches(1.7),Inches(5.4),Inches(0.4),'API 调用成本',18,B,True)
for i,(la,va,cl) in enumerate([('元数据提取 98 次','约 0.3 元',G),('查询解析每次','约 0.003 元',G),
    ('答案生成每次','约 0.005 元',G),('文本嵌入（本地 GPU）','0 元',A),('建库总成本','小于 0.5 元',D),('单次查询','约 0.01 元',O)]):
    y = Inches(2.2)+Inches(0.55)*i; tx(sl,Inches(7.2),y,Inches(3.0),Inches(0.35),la,12,GR)
    tx(sl,Inches(10.2),y,Inches(2.5),Inches(0.35),va,14,cl,True,PP_ALIGN.RIGHT)
for i,(n,l,c) in enumerate([('98','文档',B),('130+','文本块',T),('<0.5 元','总成本',G),
    ('1024','维度',A),('7 类','分类',O),('45+','测试',P)]):
    x = Inches(0.6)+Inches(2.1)*(i%3); y = Inches(6.0)+Inches(0.9)*(i//3)
    re(sl,x,y,Inches(1.8),Inches(0.8),c)
    tx(sl,x,y+Inches(0.1),Inches(1.8),Inches(0.4),n,24,WH,True,PP_ALIGN.CENTER)
    tx(sl,x,y+Inches(0.45),Inches(1.8),Inches(0.3),l,11,WH,False,PP_ALIGN.CENTER)

# ====================================================================
# 第 14 页 · 创新亮点
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,WH)
hdr(sl,'七、项目创新亮点')
hl=[('查询意图解析','LLM 将自然语言自动转为\n结构化搜索参数，用户无需\n学习搜索语法。',B),
    ('多层降级策略','解析失败 -> 全文搜索\nWhere 失败 -> 纯语义搜索\nJSON 失败 -> 正则兜底',T),
    ('Front-Matter 优先','人工 YAML 标注自动覆盖\nLLM 提取结果，尊重人工\n标注准确性。',G),
    ('来源强制追溯','System Prompt 约束加后端\n自动检查补全，100%\n答案可核查来源。',A),
    ('GPU 加速本地嵌入','Qwen3 自动 CUDA 加速\n编码速度提升 20 倍\n零 API 费用。',O),
    ('安全删除机制','delete_collection 需显式\n确认 confirm=True，防止\n误删数据。',P)]
for i,(ti,de,co) in enumerate(hl):
    x = Inches(0.6)+Inches(6.2)*(i%2); y = Inches(1.4)+Inches(1.85)*(i//2)
    re(sl,x,y,Inches(5.9),Inches(1.6),CA); re(sl,x,y,Inches(0.06),Inches(1.6),co)
    tx(sl,x+Inches(0.25),y+Inches(0.15),Inches(5.4),Inches(0.4),ti,17,co,True)
    tx(sl,x+Inches(0.25),y+Inches(0.6),Inches(5.4),Inches(0.8),de,13,TX)

# ====================================================================
# 第 15 页 · 技术选型
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,WH)
hdr(sl,'七、技术选型与权衡')
th=['技术层','选择','替代方案','选择理由']
tcx=[Inches(0.6),Inches(2.8),Inches(5.2),Inches(7.6)]
tcw=[Inches(2.1),Inches(2.3),Inches(2.3),Inches(4.0)]
for i,(x,w,h) in enumerate(zip(tcx,tcw,th)):
    re(sl,x,Inches(1.5),w,Inches(0.5),D)
    tx(sl,x+Inches(0.1),Inches(1.55),w-Inches(0.2),Inches(0.4),h,13,WH,True,PP_ALIGN.CENTER)
td=[['向量数据库','ChromaDB','Milvus/Qdrant','零部署 pip install 即用'],
    ['Embedding','Qwen3-0.6B','MiniLM-L6-v2','0.6B 参数 1024 维中文最优'],
    ['LLM','DeepSeek V4 Flash','GPT-4o-mini','0.14 美元/百万 token 性价比高'],
    ['Web 框架','Streamlit','Flask/FastAPI','纯 Python 零前端代码'],
    ['ETL 处理','Python 原生','PySpark','文档少于 100 原生 Python 足够'],
    ['PDF 解析','PyMuPDF','pdfplumber','C 底层速度快中文支持好']]
for ri,row in enumerate(td):
    y = Inches(2.0)+Inches(0.7)*ri; bg2 = CA if ri%2==0 else WH
    for ci,(x,w) in enumerate(zip(tcx,tcw)):
        re(sl,x,y,w,Inches(0.65),bg2)
        tx(sl,x+Inches(0.08),y+Inches(0.08),w-Inches(0.16),Inches(0.5),
           row[ci],11,D if ci==0 else TX,ci==0)

# ====================================================================
# 第 16 页 · 局限与改进
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,WH)
hdr(sl,'八、局限与改进方向')
for i,(li,im,so) in enumerate([('PDF 仅提取文字','含图表的文档信息丢失','增加表格和图片提取模块'),
    ('无 Rerank 重排序','Top-K 中可能含弱相关片段','接入 Cross-Encoder Reranker'),
    ('无权限隔离','所有用户共享知识库','Collection 级命名空间隔离'),
    ('分块粒度固定','不同类型文档用同一参数','按文档类型动态调整策略'),
    ('对话不跨会话','关闭后历史丢失','接入 SQLite 或 Redis 持久化'),
    ('无评测体系','无法量化检索质量','构建测试集加 Recall@K 和 MRR'),
]):
    y = Inches(1.4)+Inches(0.85)*i; bg2 = CA if i%2==0 else WH
    re(sl,Inches(0.6),y,Inches(12.1),Inches(0.75),bg2)
    re(sl,Inches(0.6),y,Inches(0.06),Inches(0.75),O)
    tx(sl,Inches(0.9),y+Inches(0.08),Inches(2.5),Inches(0.5),li,14,O,True)
    tx(sl,Inches(3.5),y+Inches(0.08),Inches(4.0),Inches(0.5),im,12,TX)
    tx(sl,Inches(7.5),y+Inches(0.08),Inches(5.0),Inches(0.5),'解决方案：'+so,12,G)

# ====================================================================
# 第 17 页 · 演示
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,WH)
hdr(sl,'九、演示方案与答辩准备')
card(sl,Inches(0.6),Inches(1.5),Inches(5.8),Inches(4.2),'CLI 演示（3 分钟）',
   ['* 知识库问答：查询提交截止日期','* 跨文档综合：什么是 RAG','* 拒答演示：查询无关问题',
    '* 展示查询解析过程'],B)
card(sl,Inches(6.9),Inches(1.5),Inches(5.8),Inches(4.2),'Streamlit 演示（5 分钟）',
   ['* 侧边栏展示文档块数和来源','* 连续提问展示对话历史','* 调试模式展示搜索词和过滤条件',
    '* 调节 Top-K 和最大距离阈值','* 展开检索来源查看分数'],T)
re(sl,Inches(0.6),Inches(5.5),Inches(12.1),Inches(1.5),CA)
re(sl,Inches(0.6),Inches(5.5),Inches(0.06),Inches(1.5),D)
tx(sl,Inches(0.9),Inches(5.6),Inches(11.5),Inches(0.3),'答辩 Q&A 准备',16,D,True)
qa=[('问：','Embedding API 限流怎么办？','答：','本地 GPU 嵌入，完全离线不存在限流问题'),
    ('问：','向量数据库崩溃怎么恢复？','答：','备份 vector_store 目录，5 分钟重建'),
    ('问：','为什么选 ChromaDB 不是 Milvus？','答：','零部署，教育场景够用；做了抽象封装可迁移'),
    ('问：','1000 篇文档怎么增量更新？','答：','当前 upsert 去重，后续可做增量索引')]
for i,(ql,q,al,a) in enumerate(qa):
    y = Inches(5.95)+Inches(0.3)*i
    tx(sl,Inches(0.9),y,Inches(0.3),Inches(0.25),ql,11,O,True)
    tx(sl,Inches(1.15),y,Inches(3.5),Inches(0.25),q,11,TX)
    tx(sl,Inches(4.65),y,Inches(0.3),Inches(0.25),al,11,G,True)
    tx(sl,Inches(4.9),y,Inches(7.5),Inches(0.25),a,11,TX)

# ====================================================================
# 第 18 页 · 致谢
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,D)
ov(sl,Inches(-0.8),Inches(-0.8),Inches(2.5),B)
ov(sl,Inches(11.0),Inches(-0.5),Inches(3.5),T)
ov(sl,Inches(10.0),Inches(5.0),Inches(5.0),B)
ov(sl,Inches(-1.5),Inches(5.5),Inches(3.0),T)
re(sl,Inches(2.5),Inches(2.5),Inches(8.3),Inches(2.5),B)
tx(sl,Inches(0.5),Inches(2.7),Inches(12.3),Inches(0.9),'谢谢聆听',44,WH,True,PP_ALIGN.CENTER)
tx(sl,Inches(0.5),Inches(3.6),Inches(12.3),Inches(0.6),'大数据学期项目计划二 -- RAG 检索增强生成助手',22,WH,False,PP_ALIGN.CENTER)
tx(sl,Inches(0.5),Inches(5.3),Inches(12.3),Inches(0.5),'代码文档齐全  |  45 个以上测试用例  |  全流程可复现',15,LG,False,PP_ALIGN.CENTER)

# ===== 保存 =====
out = Path(__file__).resolve().parent/'BigData_RAG_答辩演示.pptx'
prs.save(str(out)); print('已保存：'+str(out))
